# upload_and_extract_streamlit.py

import streamlit as st
import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation

# Function to extract text
def extract_text_from_file(file_path, file_type):
    if file_type == "pdf":
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    elif file_type == "docx":
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    elif file_type == "pptx":
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    elif file_type == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    else:
        return ""

# --- Streamlit App starts here ---
st.set_page_config(page_title="📂 Upload and Extract", layout="centered")

st.title("📚 Upload Your Lecture Notes")
st.write("👉 Please upload a `.txt`, `.pdf`, `.docx`, `.pptx`, or `.xlsx` file to extract lecture text.")

uploaded_file = st.file_uploader("Choose your lecture file:", type=["txt", "pdf", "docx", "pptx", "xlsx"])

if uploaded_file is None:
    st.info("👆 Waiting for you to upload a file...")
else:
    file_type = uploaded_file.name.split(".")[-1].lower()

    # Save file temporarily
    with open(f"temp_uploaded.{file_type}", "wb") as f:
        f.write(uploaded_file.getbuffer())

    st.success(f"✅ Uploaded: {uploaded_file.name}")

    if st.button("📖 Extract Text"):
        extracted_text = extract_text_from_file(f"temp_uploaded.{file_type}", file_type)

        if not extracted_text.strip():
            st.error("⚠️ No readable text found. Try a different file!")
        else:
            st.subheader("📝 Extracted Lecture Text")
            st.text_area(label="Extracted Text", value=extracted_text, height=400)

            st.session_state["extracted_text"] = extracted_text

# Prepare extracted text for Gemini or uAgent
if "extracted_text" in st.session_state:
    if st.button("🚀 Proceed to Summarization / Agent"):
        st.success("✅ You can now send this text to Gemini or uAgent!")
